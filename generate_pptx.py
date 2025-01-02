from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw
import os
import time 
import google.generativeai as genai

genai.configure(api_key="<Your API Key here>")


def summarize_text(text):
    model = genai.GenerativeModel("gemini-1.5-flash")
    prompt = f"""
        You have been given a very unstructed text which starts with the name of a person, the url of ther person's website, and text scraped from the website. 
        Your job is to create a short description of the person based on the given text.
        Some examples can be: 
        Eg1 : Prof and Head, Department of Urology :: LTMMC, Mumbai
        Eg2 : Consultant Urologist &  Kidney Transplant Surgeon :: Aurangabad
        Eg3 : Consultant Urologist :: MPUH, Nadiad
        Try to have the location in the end. Start the location after `::`.
        If the given given text does not have enough information to create a short description, then just return <FILL-UP>.
        Here is the text:
        \n\n{text}
    """
    response = model.generate_content(prompt)
    return response.text.strip()



def create_circular_image(input_image_path, output_image_path):
    """Create a circular version of an image."""
    img = Image.open(input_image_path).convert("RGBA")
    
    # Create a circular mask
    mask = Image.new("L", img.size, 0)
    draw = ImageDraw.Draw(mask)
    draw.ellipse((0, 0, img.size[0], img.size[1]), fill=255)

    # Apply the mask to the image
    circular_img = Image.new("RGBA", img.size)
    circular_img.paste(img, (0, 0), mask)
    circular_img.save(output_image_path)

def create_presentation(images, names, descriptions, common_text, output_file):
    # Create a PowerPoint presentation
    prs = Presentation()

    for img_path, name, description in zip(images, names, descriptions):
        # Prepare circular image
        circular_img_path = "circular_" + os.path.basename(img_path)
        create_circular_image(img_path, circular_img_path)

        # Add a blank slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Set slide background to black
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)

        # Add the circular image in the center top
        slide.shapes.add_picture(circular_img_path, Inches(3.5), Inches(1), width=Inches(3), height=Inches(3))

        # Add the name below the image
        name_box = slide.shapes.add_textbox(Inches(2.7), Inches(4), Inches(5), Inches(0.8))
        name_frame = name_box.text_frame
        name_frame.text = name
        name_frame.paragraphs[0].font.bold = True
        name_frame.paragraphs[0].font.size = Pt(32)
        name_frame.paragraphs[0].font.color.rgb = RGBColor(255, 191, 0)  # Yellow color
        name_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Add the description on the left
        desc_box = slide.shapes.add_textbox(Inches(2.7), Inches(6), Inches(5), Inches(0.8))
        desc_frame = desc_box.text_frame
        desc_frame.text = description
        desc_frame.word_wrap = True 
        desc_frame.paragraphs[0].font.size = Pt(18)
        desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        desc_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White color

        # Remove the temporary circular image
        os.remove(circular_img_path)

    # Save the presentation
    prs.save(output_file)
    

search_results_dir = "selenium_search_results"
search_results_subdirs = [name for name in os.listdir(search_results_dir) if os.path.isdir(os.path.join(search_results_dir, name))]

data = []
image_paths = []
names = []
for subdir in search_results_subdirs:
    data_path_tries = [os.path.join(search_results_dir, subdir, f"data_{i}.html") for i in range(3)]
    image_path = os.path.join(search_results_dir, subdir, "image_0.png")
    for i in range(3):
        try:
            data_content = "".join(open(data_path_tries[i], "r", encoding='utf-8').readlines())
            break
        except FileNotFoundError:
            data_content= None
            continue 
    if data_content is None:
        raise Exception("No scene")
    data.append(summarize_text(data_content))
    print(data[-1])
    image_paths.append(image_path)
    name = subdir.replace("_", " ")[:-len("Urologist")].strip()
    names.append(name)
    print(f"{name=}")
    print(f"{image_path=}")
    print(f"{data[-1]=}")
    print("")
    time.sleep(5)



output_file = "output_presentation.pptx"

create_presentation(image_paths, names, data, "", output_file)